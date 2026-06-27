#!/usr/bin/env python3
"""生成答辩 PPT (.pptx)，包含结构化图表。"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 配色 ──────────────────────────────────────────
BG    = RGBColor(0x0A, 0x10, 0x1F)
CYAN  = RGBColor(0x06, 0xB6, 0xD4)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE  = RGBColor(0xF4, 0x3F, 0x5E)
PURPLE= RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY  = RGBColor(0x94, 0xA3, 0xB8)
DARK  = RGBColor(0x64, 0x74, 0x8B)
SURFACE=RGBColor(0x13, 0x1C, 0x2E)
BORDER=RGBColor(0x1E, 0x30, 0x50)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # blank

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_box(slide, left, top, width, height, fill_color=SURFACE, border_color=BORDER, border_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    return shape

def add_rect(slide, left, top, width, height, fill_color=SURFACE, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=BORDER, width=1.5):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))  # straight
    connector.line.color.rgb = color
    connector.line.width = Pt(width)

def add_line(slide, x1, y1, x2, y2, color=BORDER, width=1.5, dashed=False):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)

def add_slide_num(slide, num, total=12):
    add_textbox(slide, 12.0, 7.05, 1.0, 0.35, f"{num}/{total}", 10, DARK, alignment=PP_ALIGN.RIGHT)

def add_tag(slide, left, top, text, bg_color, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(len(text)*0.12+0.2), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return shape

def make_slide():
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    return slide

# ═══════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ═══════════════════════════════════════════════════════
s = make_slide()
add_tag(s, 6.0, 2.0, "期末实训展示", SURFACE, CYAN)
add_textbox(s, 0, 2.6, W/Inches(1), 1.2, "基于证据约束的\n智能实习投递 Agent", 42, WHITE, True, PP_ALIGN.CENTER)
add_textbox(s, 0, 4.0, W/Inches(1), 0.6, "Evidence-Constrained Internship Application Agent", 18, GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(s, 0, 5.2, W/Inches(1), 0.4, "答辩人：XXX    |    导师：XXX    |    2026.06", 14, DARK, alignment=PP_ALIGN.CENTER)
add_slide_num(s, 1)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — 项目定位
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.4, 11, 0.6, "项目定位", 30, WHITE, True)
add_textbox(s, 0.8, 1.0, 11, 0.4, "用户做了什么 → 系统输出什么", 14, GRAY)

# 左边：用户做了什么
add_box(s, 1.0, 1.8, 4.5, 3.0, fill_color=SURFACE, border_color=GREEN)
add_textbox(s, 1.3, 1.9, 4.0, 0.4, "👤 用户做了什么", 16, GREEN, True)
items_u = ["📋 粘贴一段 JD 文本到对话框", "📄 上传 PDF 简历", "💻 丢一个 GitHub 仓库链接"]
for i, item in enumerate(items_u):
    add_textbox(s, 1.5, 2.5 + i*0.6, 3.8, 0.4, item, 14, WHITE)

# 箭头
add_textbox(s, 5.8, 2.8, 1.0, 0.5, "→", 36, CYAN, True, PP_ALIGN.CENTER)

# 右边：系统输出什么
add_box(s, 7.0, 1.8, 5.2, 3.0, fill_color=SURFACE, border_color=CYAN)
add_textbox(s, 7.3, 1.9, 4.8, 0.4, "🤖 系统输出什么", 16, CYAN, True)
items_o = ["✅ 哪些项目跟这个岗位相关？匹配度？", "✅ 简历上这些项目应该怎么写？（三档分类）", "✅ 每句话的证据来源是什么？"]
for i, item in enumerate(items_o):
    add_textbox(s, 7.5, 2.5 + i*0.6, 4.5, 0.4, item, 14, WHITE)

add_textbox(s, 0, 5.4, W/Inches(1), 0.5, "核心原则：所有输出必须基于证据，不得编造", 18, CYAN, True, PP_ALIGN.CENTER)
add_slide_num(s, 2)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — Demo：用户做了什么
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.4, 11, 0.6, "Demo：用户做了什么", 30, WHITE, True)
add_textbox(s, 0.8, 1.0, 11, 0.4, "真实场景 — 用户粘贴 JD、上传 PDF 简历、丢 GitHub 链接", 14, GRAY)

# JD box
add_box(s, 0.8, 1.7, 6.5, 4.5, fill_color=SURFACE, border_color=AMBER)
add_textbox(s, 1.1, 1.8, 6.0, 0.35, "📋 粘贴的 JD（纯文本）", 15, AMBER, True)
jd_lines = [
    "Agent开发实习生 (140-190元/天) · 深圳 · 5天/周",
    "",
    "岗位职责：",
    "· 负责 AI Agent 系统的开发和优化",
    "· 记忆模块、规划模块和工具调用模块",
    "· 开发基于大语言模型(LLM)的智能体系统",
    "· 设计 Agent 的工具使用接口",
    "· 开发 Agent 的知识库管理系统",
    "",
    "任职要求：",
    "· 熟悉 Python、Java · 熟悉 LLM 原理和应用",
    "· 了解机器学习、深度学习等 AI 技术",
]
for i, line in enumerate(jd_lines):
    c = WHITE if line.startswith("·") or line.startswith("Agent") else (AMBER if "：" in line else GRAY)
    add_textbox(s, 1.3, 2.3 + i*0.3, 5.8, 0.28, line, 12, c)

# PDF + GitHub boxes
add_box(s, 7.8, 1.7, 4.7, 2.0, fill_color=SURFACE, border_color=GREEN)
add_textbox(s, 8.1, 1.8, 4.2, 0.35, "📄 上传的 PDF 简历", 15, GREEN, True)
add_textbox(s, 8.3, 2.3, 4.0, 1.2, "陈小明 · XX大学 化学+CS\nPython(3年) · FastAPI · Docker\n4个项目 · 2段实习经历", 13, WHITE)

add_box(s, 7.8, 4.0, 4.7, 2.0, fill_color=SURFACE, border_color=PURPLE)
add_textbox(s, 8.1, 4.1, 4.2, 0.35, "💻 丢的 GitHub 链接", 15, PURPLE, True)
add_textbox(s, 8.3, 4.6, 4.0, 1.2, "github.com/yfn-1116/\ncareer-agent-assistant\nRAG 求职 Agent · 15 ⭐", 13, WHITE)

add_textbox(s, 0, 6.6, W/Inches(1), 0.4, "用户就做了这三件事。LLM 自己判断意图，自动调 Agent 处理。", 14, CYAN, alignment=PP_ALIGN.CENTER)
add_slide_num(s, 3)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — Demo：哪些项目相关？
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.4, 11, 0.6, "Demo：哪些项目跟这个 JD 相关？", 30, WHITE, True)
add_textbox(s, 0.8, 1.0, 11, 0.4, "系统从你的经历中检索并匹配", 14, GRAY)

# 3 evidence cards
ev = [
    ("✅ SmartApplyAgent — RAG 智能求职助手", "implemented · 匹配 0.89", GREEN,
     "匹配 JD 需求：Agent系统开发 · 工具调用模块 · 知识库管理 · 对话系统 · Python\n来源: github.com/yfn-1116/career-agent-assistant"),
    ("✅ 智能滴定系统 — 视觉滴定终点检测", "implemented · 匹配 0.62", GREEN,
     "匹配 JD 需求：Python · 工程实践 · 系统设计能力\n来源: github.com/yfn-1116/auto-titration-system"),
    ("⚠️ PolyU 智能导航 — 校园路径规划", "designed · 匹配 0.55", AMBER,
     "匹配 JD 需求：系统架构设计 · API 集成\n来源: github.com/yfn-1116/polyu-smart-journey"),
]
for i, (title, badge, color, desc) in enumerate(ev):
    y = 1.7 + i * 1.55
    add_box(s, 0.8, y, 11.5, 1.35, fill_color=SURFACE, border_color=color)
    add_textbox(s, 1.1, y + 0.08, 8.0, 0.35, title, 16, WHITE, True)
    add_tag(s, 9.8, y + 0.1, badge, SURFACE, color)
    add_textbox(s, 1.1, y + 0.55, 10.5, 0.7, desc, 12, GRAY)

# Coverage
add_textbox(s, 0.8, 6.4, 11.5, 0.5, "JD 需求覆盖：✅ Agent系统开发  ✅ 工具调用  ✅ 知识库管理  ✅ 对话系统  ✅ Python    ⚠️ Java    📋 机器学习", 13, WHITE)
add_slide_num(s, 4)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — Demo：简历怎么写
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.4, 11, 0.6, "Demo：简历怎么写？（三档分类）", 30, WHITE, True)
add_textbox(s, 0.8, 1.0, 11, 0.4, '基于证据等级，不是基于模型"想象"', 14, GRAY)

tiers = [
    ("✅ 可以直接写进简历的", "implemented 证据", GREEN,
     '"独立开发 AI Agent 智能求职助手，基于 LangGraph 实现多 Agent 协作框架。设计并实现\n工具调用模块(task_agent)，支持 Agent 动态调度与递归生成。构建 RAG 知识库管理系统\n(BM25+向量+重排序)。开发 Streamlit 对话交互界面。"\n[证据: career-agent-assistant | 等级: implemented]'),
    ("⚠️ 需要确认后才能写的", "designed 证据", AMBER,
     '"设计基于 FastAPI + PostgreSQL 的校园导航后端，实现 RESTful API 与外部地图服务集成"\n[证据: polyu-smart-journey | 等级: designed — 仅设计，需确认]'),
    ("📋 建议补充实践后再写的", "无对应证据", DARK,
     'Java — 知识库中无相关实现经历，建议补充实践    |    机器学习/深度学习 — 无直接项目经验，建议补充 Demo'),
]
for i, (title, badge, color, desc) in enumerate(tiers):
    y = 1.7 + i * 1.7
    add_box(s, 0.8, y, 11.5, 1.45, fill_color=SURFACE, border_color=color)
    add_textbox(s, 1.1, y + 0.08, 6.5, 0.35, title, 16, WHITE, True)
    add_tag(s, 7.8, y + 0.1, badge, SURFACE, color)
    add_textbox(s, 1.1, y + 0.55, 10.8, 0.85, desc, 11, GRAY)

add_textbox(s, 0, 6.8, W/Inches(1), 0.4, "核心：基于真实证据，告诉你能写什么、不能写什么、还需要补什么", 15, CYAN, True, PP_ALIGN.CENTER)
add_slide_num(s, 5)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — Agent 四要素
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.3, 11, 0.6, "Agent 系统的四个要素", 30, WHITE, True)
add_textbox(s, 0.8, 0.85, 11, 0.4, "Demo 背后是这个四层架构在运转", 14, GRAY)

# 4 boxes in 2x2 grid
boxes = [
    (0.6, 1.6, "🔍 Perception 感知", GREEN, ["理解用户输入", "JD文本→结构化技能", "PDF→提取文本", "GitHub链接→仓库内容"]),
    (6.8, 1.6, "🧠 Planning 规划", AMBER, ["决定做什么", "LLM识别意图", "决定调哪个Agent", "决定并行/串行"]),
    (0.6, 4.2, "⚡ Action 执行", CYAN, ["执行操作", "RAG检索知识库", "匹配分析", "证据校验"]),
    (6.8, 4.2, "💾 Memory 记忆", PURPLE, ["记住一切", "短期:对话上下文", "长期:JSONL持久化", "知识库:RAG索引"]),
]
for x, y, title, color, items in boxes:
    add_box(s, x, y, 5.5, 2.3, fill_color=SURFACE, border_color=color)
    add_textbox(s, x + 0.3, y + 0.15, 4.8, 0.4, title, 18, color, True)
    for j, item in enumerate(items):
        add_textbox(s, x + 0.3, y + 0.7 + j*0.35, 4.8, 0.3, item, 13, GRAY if j > 0 else WHITE)

# Cross-layer RAG annotation
add_textbox(s, 0, 6.7, W/Inches(1), 0.4, "═══ RAG 横切四层 ═══   写入: Perception→Memory  |  读取: Memory→Action", 12, GRAY, alignment=PP_ALIGN.CENTER)
add_slide_num(s, 6)

# ═══════════════════════════════════════════════════════
# SLIDE 7 — 系统总体架构
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.25, 11, 0.6, "系统总体架构", 30, WHITE, True)
add_textbox(s, 0.8, 0.75, 11, 0.35, "用户操作 → 调度层决策 → 能力层执行 → 基础设施 → 存储", 14, GRAY)

layers = [
    (0.8, 1.3, 11.5, 0.7, "📋 粘贴JD    📄 上传PDF    💻 丢GitHub链接", SURFACE, BORDER, WHITE, 14),
    (0.8, 2.2, 11.5, 1.5, "调度层：OrchestratorAgent", SURFACE, CYAN, CYAN, 16),
    (0.8, 3.9, 3.6, 1.4, "能力层\n16个Tool\nparse_jd→retrieve→\nanalyze→generate→check", SURFACE, GREEN, WHITE, 12),
    (4.7, 3.9, 3.6, 1.4, "基础设施层\nLLM(Qwen/DeepSeek/Mock)\nEmbedding(1024维)\nRAG Pipeline · MCP", SURFACE, AMBER, WHITE, 12),
    (8.6, 3.9, 3.7, 1.4, "存储层\n短期(20条)\n长期(JSONL+BM25)\n知识库(chunks.jsonl)", SURFACE, PURPLE, WHITE, 12),
]
for x, y, w, h, text, fill, border, tcolor, fsize in layers:
    if "调度层" in str(text):
        add_box(s, x, y, w, h, fill_color=fill, border_color=border)
        add_textbox(s, x + 0.3, y + 0.1, w - 0.5, 0.35, "调度层：OrchestratorAgent", 16, CYAN, True)
        add_textbox(s, x + 0.3, y + 0.55, w - 0.5, 0.85,
            'LLM 收到消息 → 判断意图："这是岗位分析任务"\n→ 决定：先解析JD+索引知识库 → 再检索匹配 → 再分析 → 最后生成\n不是写死的 if-else，是 LLM 自己做的决策。LLM 不可用时回退 LangGraph DAG。', 12, GRAY)
    elif "粘贴" in str(text):
        add_textbox(s, x, y, w, h, text, fsize, tcolor, alignment=PP_ALIGN.CENTER)
    else:
        add_box(s, x, y, w, h, fill_color=fill, border_color=border)
        lines = text.split('\n')
        for j, line in enumerate(lines):
            c = tcolor if j == 0 else GRAY
            sz = 14 if j == 0 else 12
            add_textbox(s, x + 0.2, y + 0.1 + j*0.28, w - 0.4, 0.3, line, sz, c, j == 0)

add_textbox(s, 0, 5.6, W/Inches(1), 0.4, "关键设计：LLM 不自己干活 —— 它生成子 Agent 去干活。接下来展开 →", 14, CYAN, True, PP_ALIGN.CENTER)
add_slide_num(s, 7)

# ═══════════════════════════════════════════════════════
# SLIDE 8 — 主 Agent 决策过程
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.25, 11, 0.6, "主 Agent 怎么分派任务？", 30, WHITE, True)
add_textbox(s, 0.8, 0.75, 11, 0.35, "OrchestratorAgent 的 LLM 决策循环 — 以这个 JD 为例", 14, GRAY)

steps = [
    (1, "LLM判断→需要解析JD并索引知识库", '调 task_agent("解析JD+索引PDF和GitHub", [parse_jd, web_search]) → SubAgent-1', False, GREEN),
    (2, "LLM判断→需要检索匹配的经历  ⚡与Step1并行", '调 task_agent("检索Agent开发相关经历", [retrieve_profile, github_repo]) → SubAgent-2', True, GREEN),
    (3, "LLM判断→需要分析匹配度  🔗依赖Step2结果", '调 task_agent("分析JD和证据的匹配度", [analyze_match, grade_retrieval]) → SubAgent-3', False, AMBER),
    (4, "LLM判断→生成简历建议", '调 task_agent("生成简历建议", [generate_answer, check_faithfulness, write_report]) → SubAgent-4', False, AMBER),
    (5, "LLM判断→任务完成，回复用户", "不再调 Tool，直接输出最终答案给用户", False, WHITE),
]
for i, (num, title, detail, is_parallel, color) in enumerate(steps):
    y = 1.4 + i * 1.05
    add_box(s, 1.5, y, 10.0, 0.85, fill_color=SURFACE, border_color=color)
    add_textbox(s, 1.1, y + 0.1, 0.5, 0.3, str(num), 14, CYAN, True, PP_ALIGN.CENTER)
    add_textbox(s, 1.8, y + 0.08, 9.5, 0.35, title, 14, WHITE, True)
    add_textbox(s, 1.8, y + 0.45, 9.5, 0.35, detail, 12, GRAY)
    if is_parallel:
        add_tag(s, 10.0, y + 0.05, "⚡并行", SURFACE, GREEN)

add_textbox(s, 0, 6.8, W/Inches(1), 0.4, "主 Agent 不干活，只做决策。干活的是子 Agent。", 15, CYAN, True, PP_ALIGN.CENTER)
add_slide_num(s, 8)

# ═══════════════════════════════════════════════════════
# SLIDE 9 — 子 Agent 执行与回传
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.25, 11, 0.6, "子 Agent 怎么执行、怎么回传？", 30, WHITE, True)
add_textbox(s, 0.8, 0.75, 11, 0.35, "以 SubAgent-1（感知·解析专家）为例", 14, GRAY)

# Left: SubAgent-1 detail
add_box(s, 0.6, 1.35, 6.5, 5.0, fill_color=SURFACE, border_color=CYAN)
add_textbox(s, 0.9, 1.45, 5.8, 0.35, "📄 SubAgent-1：感知·解析专家", 16, CYAN, True)
add_tag(s, 0.9, 1.9, "Role: 信息采集专家", SURFACE, GREEN)
add_tag(s, 3.2, 1.9, "Session: 独立上下文", SURFACE, CYAN)
add_tag(s, 5.1, 1.9, "max_steps: 8", SURFACE, AMBER)
add_textbox(s, 0.9, 2.35, 5.8, 0.6, 'System Prompt: "你是信息采集专家，解析所有输入并索引到知识库。"\nPermissions: [parse_jd] [web_search]', 12, GRAY)
add_textbox(s, 0.9, 3.05, 5.8, 0.25, "执行过程：", 13, WHITE, True)
exec_steps = [
    "① 调 parse_jd(JD文本) → {job_title, hard_skills: [Python, Java, LLM, Agent, 工具调用, 知识库...]}",
    "② 调 FileLoader(PDF) → pypdf 提取文本",
    "③ TextChunker 分块(800字) → SHA256 ID",
    "④ 写入 chunks.jsonl → 3个chunk索引完成",
]
for i, es in enumerate(exec_steps):
    add_textbox(s, 0.9, 3.35 + i*0.35, 5.8, 0.3, es, 11, GRAY)

# Right: communication flow
add_box(s, 7.5, 1.35, 5.0, 5.0, fill_color=SURFACE, border_color=AMBER)
add_textbox(s, 7.8, 1.45, 4.4, 0.35, "通信流程", 14, AMBER, True)

add_box(s, 7.9, 2.0, 4.2, 1.1, fill_color=SURFACE, border_color=CYAN)
add_textbox(s, 8.1, 2.05, 3.8, 0.8, '① 主Agent调 task_agent\ntask="解析JD文本，索引PDF\n和GitHub到知识库"\nallowed_tools=[parse_jd, web_search]', 11, WHITE)

add_textbox(s, 9.5, 3.2, 1.0, 0.3, "▼", 14, CYAN, True, PP_ALIGN.CENTER)

add_box(s, 7.9, 3.5, 4.2, 0.9, fill_color=SURFACE, border_color=GREEN)
add_textbox(s, 8.1, 3.55, 3.8, 0.8, "② SubAgent-1 独立执行\n自己的System Prompt+上下文\n不受其他Agent干扰", 11, WHITE)

add_textbox(s, 9.5, 4.5, 1.0, 0.3, "▼", 14, CYAN, True, PP_ALIGN.CENTER)

add_box(s, 7.9, 4.8, 4.2, 0.9, fill_color=SURFACE, border_color=PURPLE)
add_textbox(s, 8.1, 4.85, 3.8, 0.8, '③ 返回 SubAgentResult\n{answer:"JD已解析,PDF索引\n3chunk", tools_called:["parse_jd"]}', 11, WHITE)

add_textbox(s, 7.8, 5.85, 4.4, 0.35, "子Agent也能再调task_agent→孙Agent(递归)", 11, GRAY, alignment=PP_ALIGN.CENTER)
add_slide_num(s, 9)

# ═══════════════════════════════════════════════════════
# SLIDE 10 — 四个子 Agent 全貌
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.2, 11, 0.55, "四个子 Agent 全貌", 30, WHITE, True)
add_textbox(s, 0.8, 0.7, 11, 0.3, "Role（角色）× Permissions（权限）× Session（上下文）三维解耦", 14, GRAY)

# 4 agent cards
cards = [
    (0.3, 1.3, "📄 SubAgent-1", "感知·解析专家", ["parse_jd", "web_search"], "⚡ 并行", GREEN,
     "← JD+PDF+GitHub\n→ ParsedJD + 索引确认"),
    (3.2, 1.3, "🔍 SubAgent-2", "检索·知识专家", ["retrieve_profile", "github_repo"], "⚡ 并行", GREEN,
     "← JD技能列表\n→ Evidence[3] + Grade"),
    (6.1, 1.3, "⚖️ SubAgent-3", "分析·匹配专家", ["analyze_match", "grade_retrieval"], "🔗 串行", AMBER,
     "← ParsedJD + Evidence[]\n→ MatchAnalysis"),
    (9.0, 1.3, "✏️ SubAgent-4", "生成·输出专家", ["generate_answer", "write_report"], "🔗 串行", AMBER,
     "← MatchAnalysis\n→ bullets+message+report"),
]
for x, y, name, role, perms, exec_mode, color, io in cards:
    add_box(s, x, y, 2.7, 3.5, fill_color=SURFACE, border_color=color)
    add_textbox(s, x + 0.15, y + 0.1, 2.3, 0.3, name, 14, WHITE, True)
    add_textbox(s, x + 0.15, y + 0.4, 2.3, 0.2, role, 10, GRAY)
    add_textbox(s, x + 0.15, y + 0.7, 2.3, 0.2, f"Permissions:", 9, DARK)
    for j, p in enumerate(perms):
        add_tag(s, x + 0.15 + j*1.2, y + 0.9, p, SURFACE, color)
    add_textbox(s, x + 0.15, y + 1.3, 2.3, 0.25, f"Session: 独立上下文", 10, GRAY)
    add_textbox(s, x + 0.15, y + 1.6, 2.3, 0.25, f"{exec_mode}", 12, color, True)
    add_textbox(s, x + 0.15, y + 1.95, 2.3, 0.25, "做什么:", 10, DARK)
    add_textbox(s, x + 0.15, y + 2.2, 2.3, 0.8, io, 10, GRAY)

# Recursive branch for SubAgent-3
add_box(s, 6.3, 5.1, 2.4, 1.5, fill_color=SURFACE, border_color=PURPLE)
add_textbox(s, 6.45, 5.15, 2.0, 0.25, "↳ SubAgent-3-1 (递归!)", 11, PURPLE, True)
add_textbox(s, 6.45, 5.45, 2.0, 0.2, "✅ 证据校验专员", 10, GRAY)
add_tag(s, 6.45, 5.7, "check_faithfulness", SURFACE, PURPLE)
add_textbox(s, 6.45, 6.05, 2.0, 0.5, "→ can_write:2\n  confirm:1\n  blocked:0", 10, GRAY)

# Parallel/Serial/Recursive legend
add_textbox(s, 0.3, 6.8, 12.0, 0.5,
    "● 并行：SubAgent-1 ∥ SubAgent-2（无数据依赖，LLM同步生成）    ● 串行：SubAgent-3等SubAgent-2，SubAgent-4等SubAgent-3    ● 递归：SubAgent-3→SubAgent-3-1",
    11, GRAY)
add_slide_num(s, 10)

# ═══════════════════════════════════════════════════════
# SLIDE 11 — RAG Pipeline
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0.8, 0.2, 11, 0.55, "RAG：知识库的数据怎么进去、怎么查出来", 28, WHITE, True)
add_textbox(s, 0.8, 0.65, 11, 0.3, "写入路径（Perception→Memory） + 读取路径（Memory→Action）", 13, GRAY)

# Left: Write path
add_box(s, 0.4, 1.2, 5.8, 4.2, fill_color=SURFACE, border_color=CYAN)
add_textbox(s, 0.7, 1.3, 5.0, 0.3, "📥 写入路径 (Perception → Memory)", 15, CYAN, True)
write_steps = [
    "① PDF简历 → pypdf 提取文本",
    "② GitHub链接 → MCP(JSON-RPC) 调 npx 子进程读 README",
    "   失败回退：urllib 直接 GET raw 内容",
    "③ TextChunker 滑动窗口分块",
    "   chunk_size=800, overlap=100, step=700",
    "④ 每 chunk 生成 SHA256 document_id",
    "⑤ 追加写入 chunks.jsonl",
]
for i, ws in enumerate(write_steps):
    c = GRAY if ws.startswith("  ") else WHITE
    add_textbox(s, 0.9, 1.8 + i*0.32, 4.8, 0.28, ws, 11, c)

# Right: Read path
add_box(s, 6.5, 1.2, 6.3, 4.2, fill_color=SURFACE, border_color=AMBER)
add_textbox(s, 6.8, 1.3, 5.5, 0.3, "📤 读取路径 (Memory → Action)", 15, AMBER, True)
add_textbox(s, 6.8, 1.65, 5.5, 0.25, '查询词: "Python Agent LLM 工具调用 知识库管理"', 11, GRAY)
read_steps = [
    ("①", "jieba 中文分词 → token 集合"),
    ("②", "双路并行粗召回 (各取 top 80)"),
    ("", "  BM25 关键词 ⊕ Qwen Embedding 语义(1024维)"),
    ("", "  互补：BM25 保准确，Embedding 保召回"),
    ("③", "RRF 融合 (k=60, → top 30)"),
    ("", "  只关心排序位置，不关心分数绝对值，零参数"),
    ("④", "CrossEncoder 精排 (→ top 5)"),
    ("", "  bge-reranker-base 279M · query+chunk 过 Transformer"),
    ("⑤", "FaithfulnessChecker (阈值 0.75)"),
    ("⑥", "EvidenceGate 四级约束"),
    ("", "  ✅implemented→可写  ⚠️designed→需确认  ❌planned→禁止"),
]
for i, (num, desc) in enumerate(read_steps):
    y = 1.95 + i*0.28
    if num:
        add_textbox(s, 6.8, y, 0.4, 0.25, num, 10, CYAN, True)
        add_textbox(s, 7.2, y, 4.8, 0.25, desc, 10, WHITE, True)
    else:
        add_textbox(s, 7.4, y, 4.6, 0.25, desc.strip(), 9, GRAY)

# Bottom: fallback
add_box(s, 0.4, 5.6, 12.5, 0.7, fill_color=SURFACE, border_color=AMBER)
add_textbox(s, 0.7, 5.7, 11.8, 0.5,
    "⚡ API 不可用 → 轻量回退：BM25 + 五维规则打分(技能0.35+来源0.25+密度0.20+长度0.10+去重0.10)，<50ms，零网络依赖\n代码：hybrid_retriever.py · rrf_fusion.py · cross_encoder_reranker.py · gate.py:46", 11, AMBER)
add_slide_num(s, 11)

# ═══════════════════════════════════════════════════════
# SLIDE 12 — 回顾
# ═══════════════════════════════════════════════════════
s = make_slide()
add_textbox(s, 0, 0.6, W/Inches(1), 0.6, "回顾三个核心设计", 32, WHITE, True, PP_ALIGN.CENTER)

review = [
    (0.6, 2.0, "🌳 Agent 递归树", "主Agent决策→子Agent执行\nRole×Permissions×Session\n并行+串行+递归", CYAN),
    (4.7, 2.0, "🔧 RAG 全链路", "写入: PDF/GitHub→分块→索引\n读取: BM25+Embedding→RRF→精排\nAPI不可用时自动降级", AMBER),
    (8.8, 2.0, "🛡️ EvidenceGate", "implemented→✅可写\ndesigned→⚠️需确认\nplanned→❌禁止", GREEN),
]
for x, y, title, desc, color in review:
    add_box(s, x, y, 3.8, 3.0, fill_color=SURFACE, border_color=color)
    add_textbox(s, x + 0.3, y + 0.3, 3.1, 0.4, title, 20, color, True, PP_ALIGN.CENTER)
    add_textbox(s, x + 0.3, y + 1.0, 3.1, 1.5, desc, 14, GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0, 5.4, W/Inches(1), 0.4, "三个局限：语义检索依赖外部 API · 大规模并发未验证 · 测评数据集有限", 13, GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(s, 0, 6.2, W/Inches(1), 0.5, "欢迎各位老师提问 🙋", 22, WHITE, True, PP_ALIGN.CENTER)
add_slide_num(s, 12)

# ── 保存 ──────────────────────────────────────────
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'outputs')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'defense_presentation.pptx')
prs.save(out_path)
print(f"✅ PPT saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
